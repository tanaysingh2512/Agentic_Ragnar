import os
from pptx import Presentation

def pptx_folder_to_txt(folder_path, txt_path):
    with open(txt_path, 'w', encoding='utf-8') as txt_file:
        for filename in os.listdir(folder_path):
            if filename.lower().endswith('.pptx'):
                pptx_file = os.path.join(folder_path, filename)
                prs = Presentation(pptx_file)
                txt_file.write(f"===== {filename} =====\n")
                for slide_num, slide in enumerate(prs.slides, start=1):
                    txt_file.write(f"Slide {slide_num}:\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            txt_file.write(shape.text.strip() + '\n')
                    txt_file.write('\n' + '-'*40 + '\n\n')

# Usage:
pptx_folder_to_txt("Vanliga frågor", "knowledge-base.txt")
